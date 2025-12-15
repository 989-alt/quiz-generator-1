# ==========================================
# 1. 환경 설정 및 라이브러리 설치 안내
# ==========================================
# 터미널에서 아래 명령어로 필요한 라이브러리를 설치하세요.
# pip install streamlit google-generativeai PyPDF2 python-pptx pandas

import streamlit as st
import google.generativeai as genai
import PyPDF2
from pptx import Presentation
import pandas as pd
import json
import io

# ==========================================
# 2. 유틸리티 함수 (텍스트 추출 및 API 호출)
# ==========================================

def extract_text_from_file(uploaded_file):
    """업로드된 파일에서 텍스트를 추출합니다 (PDF, PPTX, TXT)."""
    text = ""
    try:
        if uploaded_file.name.endswith('.pdf'):
            reader = PyPDF2.PdfReader(uploaded_file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif uploaded_file.name.endswith('.pptx'):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif uploaded_file.name.endswith('.txt'):
            text = uploaded_file.read().decode("utf-8")
    except Exception as e:
        st.error(f"파일 읽기 오류: {e}")
        return None
    return text

def get_gemini_response(api_key, prompt, model_name="gemini-1.5-flash"):
    """Gemini API를 호출하여 응답을 받습니다."""
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(model_name)
        # JSON 포맷을 강제하기 위한 설정 (1.5 Flash 최신 버전 지원 시)
        generation_config = genai.GenerationConfig(response_mime_type="application/json")
        response = model.generate_content(prompt, generation_config=generation_config)
        return response.text
    except Exception as e:
        st.error(f"Gemini API 호출 오류: {e}")
        return None

def parse_quiz_json(json_str):
    """API 응답(문자열)을 Python 리스트/딕셔너리로 변환합니다."""
    try:
        # 가끔 마크다운 코드 블록(```json ... ```)이 포함될 경우 제거
        clean_str = json_str.strip()
        if clean_str.startswith("```json"):
            clean_str = clean_str.replace("```json", "", 1)
        if clean_str.startswith("```"):
            clean_str = clean_str.replace("```", "", 1)
        if clean_str.endswith("```"):
            clean_str = clean_str.rsplit("```", 1)[0]
        
        return json.loads(clean_str)
    except json.JSONDecodeError:
        st.error("AI 응답을 JSON으로 파싱하는데 실패했습니다. 다시 시도해주세요.")
        return []

# ==========================================
# 3. 퀴즈 생성 프롬프트 빌더
# ==========================================

def build_prompt(context_text, num_questions=5):
    return f"""
    당신은 전문적인 교사입니다. 아래 제공된 학습 자료를 바탕으로 {num_questions}개의 객관식 퀴즈(4지 선다형)를 만들어주세요.
    
    [조건]
    1. 반드시 JSON 형식의 리스트로 반환해야 합니다.
    2. 각 문제는 다음 키를 가져야 합니다: "question", "options", "answer", "explanation".
    3. "options"는 반드시 4개의 선택지를 가진 리스트여야 합니다.
    4. "answer"는 정답인 보기를 그대로 적어주세요.
    5. 언어는 한국어여야 합니다.

    [학습 자료]
    {context_text[:10000]} (내용이 너무 길면 일부만 사용됨)

    [출력 예시]
    [
        {{
            "question": "대한민국의 수도는?",
            "options": ["부산", "서울", "대구", "인천"],
            "answer": "서울",
            "explanation": "대한민국의 수도는 서울입니다."
        }}
    ]
    """

def build_single_regen_prompt(context_text):
    """단일 문제 재생성을 위한 프롬프트"""
    return f"""
    당신은 전문적인 교사입니다. 아래 학습 자료를 바탕으로 기존과 다른 새로운 객관식 퀴즈(4지 선다형) 1개를 만들어주세요.
    반드시 리스트 형식이 아닌 단일 JSON 객체로 반환하세요.

    [학습 자료]
    {context_text[:5000]}
    
    [출력 예시]
    {{
        "question": "새로운 문제 내용...",
        "options": ["보기1", "보기2", "보기3", "보기4"],
        "answer": "보기1",
        "explanation": "해설..."
    }}
    """

# ==========================================
# 4. 콜백 함수 (상태 관리)
# ==========================================

def delete_question(index):
    """퀴즈 리스트에서 특정 인덱스의 문제를 삭제합니다."""
    if "quiz_data" in st.session_state:
        del st.session_state.quiz_data[index]
        st.toast("문제가 삭제되었습니다.", icon="🗑️")

def regenerate_question(index, api_key, context_text):
    """특정 인덱스의 문제를 새로운 문제로 교체합니다."""
    if not api_key:
        st.warning("API Key가 필요합니다.")
        return

    with st.spinner("문제 재생성 중..."):
        prompt = build_single_regen_prompt(context_text)
        res_text = get_gemini_response(api_key, prompt)
        if res_text:
            new_quiz = parse_quiz_json(res_text)
            # 단일 객체인지 리스트인지 확인 후 처리
            if isinstance(new_quiz, list) and len(new_quiz) > 0:
                 st.session_state.quiz_data[index] = new_quiz[0]
            elif isinstance(new_quiz, dict):
                 st.session_state.quiz_data[index] = new_quiz
            st.toast("문제가 재생성되었습니다.", icon="🔄")

# ==========================================
# 5. 메인 앱 UI 구조
# ==========================================

def main():
    st.set_page_config(page_title="AI 퀴즈 생성기", page_icon="📝", layout="wide")

    # --- Session State 초기화 ---
    if "quiz_data" not in st.session_state:
        st.session_state.quiz_data = []
    if "source_text" not in st.session_state:
        st.session_state.source_text = ""

    # --- 사이드바: 설정 ---
    with st.sidebar:
        st.header("⚙️ 설정")
        api_key = st.text_input("Gemini API Key", type="password", help="Google AI Studio에서 발급받은 키를 입력하세요.")
        st.info("지원 모델: gemini-1.5-flash")
        st.markdown("---")
        st.markdown("**사용 방법**")
        st.markdown("1. API Key 입력\n2. 강의 자료 업로드\n3. 퀴즈 생성 클릭\n4. 수정 및 다운로드")

    # --- 메인 영역 ---
    st.title("📝 AI 기반 수업 자료 퀴즈 생성기")
    st.markdown("PDF, PPT, 텍스트 파일을 업로드하면 **4지 선다형 퀴즈**를 자동으로 생성합니다.")

    # 파일 업로더
    uploaded_file = st.file_uploader("강의 자료 업로드 (PDF, PPTX, TXT)", type=["pdf", "pptx", "txt"])

    if uploaded_file:
        # 텍스트 추출 (최초 1회 또는 파일 변경 시)
        # 파일이 바뀌었는지 체크하는 로직을 간단하게 구현하기 위해 바로 추출
        extracted_text = extract_text_from_file(uploaded_file)
        if extracted_text:
            st.session_state.source_text = extracted_text
            with st.expander("추출된 텍스트 미리보기"):
                st.text(extracted_text[:1000] + "...")
        else:
            st.error("텍스트 추출에 실패했습니다.")

    # 퀴즈 생성 버튼
    if st.button("🚀 퀴즈 생성하기", type="primary"):
        if not api_key:
            st.warning("⚠️ 왼쪽 사이드바에 Gemini API Key를 먼저 입력해주세요.")
        elif not st.session_state.source_text:
            st.warning("⚠️ 파일을 업로드해주세요.")
        else:
            with st.spinner("AI가 내용을 분석하고 퀴즈를 만들고 있습니다..."):
                prompt = build_prompt(st.session_state.source_text, num_questions=5)
                res_text = get_gemini_response(api_key, prompt)
                
                if res_text:
                    quiz_list = parse_quiz_json(res_text)
                    if quiz_list:
                        st.session_state.quiz_data = quiz_list
                        st.success(f"{len(quiz_list)}개의 문제가 생성되었습니다!")
                    else:
                        st.error("퀴즈 데이터 파싱 실패. 다시 시도해주세요.")

    st.markdown("---")

    # --- 퀴즈 표시 및 관리 영역 ---
    if st.session_state.quiz_data:
        st.subheader("✅ 생성된 퀴즈 리스트")
        
        # 반복문을 통해 각 퀴즈 표시
        # enumerate를 사용하여 인덱스 확보 (삭제/재생성 시 필요)
        for i, q_item in enumerate(st.session_state.quiz_data):
            with st.container(border=True):
                # 상단: 문제 제목
                col_q, col_btn = st.columns([8, 2])
                with col_q:
                    st.markdown(f"**Q{i+1}. {q_item.get('question')}**")
                
                # 우측 상단: 관리 버튼 (재생성, 삭제)
                with col_btn:
                    c1, c2 = st.columns(2)
                    with c1:
                        st.button("🔄", key=f"regen_{i}", help="이 문제 재생성", 
                                  on_click=regenerate_question, 
                                  args=(i, api_key, st.session_state.source_text))
                    with c2:
                        st.button("🗑️", key=f"del_{i}", help="이 문제 삭제", 
                                  on_click=delete_question, 
                                  args=(i,))

                # 보기 표시
                options = q_item.get('options', [])
                for idx, opt in enumerate(options):
                    st.text(f"{idx+1}) {opt}")

                # 정답 및 해설 (토글)
                with st.expander("정답 및 해설 확인"):
                    st.success(f"정답: {q_item.get('answer')}")
                    st.info(f"해설: {q_item.get('explanation')}")

        # --- CSV 다운로드 영역 ---
        st.markdown("### 📥 결과 저장")
        
        # 데이터프레임 변환
        csv_data = []
        for i, q in enumerate(st.session_state.quiz_data):
            row = {
                "번호": i + 1,
                "문제": q.get("question"),
                "보기1": q.get("options")[0] if len(q.get("options")) > 0 else "",
                "보기2": q.get("options")[1] if len(q.get("options")) > 1 else "",
                "보기3": q.get("options")[2] if len(q.get("options")) > 2 else "",
                "보기4": q.get("options")[3] if len(q.get("options")) > 3 else "",
                "정답": q.get("answer"),
                "해설": q.get("explanation")
            }
            csv_data.append(row)
        
        df = pd.DataFrame(csv_data)
        csv_utf8 = df.to_csv(index=False, encoding='utf-8-sig') # 한글 깨짐 방지
        
        st.download_button(
            label="CSV 파일로 다운로드",
            data=csv_utf8,
            file_name="generated_quiz.csv",
            mime="text/csv"
        )

if __name__ == "__main__":
    main()